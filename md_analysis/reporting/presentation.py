"""
PowerPoint Presentation Generation Module.

This module contains functions for generating a PowerPoint presentation
summarizing results from multiple MD analysis runs.
"""

import os
import logging
import json
from collections import defaultdict
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

# Set up logger
logger = logging.getLogger(__name__)

def Create_PPT(unique_dirs, com_averages):
    """
    Creates a PowerPoint presentation summarizing results from multiple runs.
    Includes raw/filtered distance plots, orientation plots, and a COM summary table.
    Clearly identifies control systems vs toxin-channel systems.

    Args:
        unique_dirs (list): List of paths to the run directories to include.
        com_averages (dict): Dictionary mapping system names to lists of filtered
                             average COM distances for the runs within that system.
    """
    logger.info(f"Creating PowerPoint summary for {len(unique_dirs)} runs...")
    prs = Presentation()

    # Group runs by system (assuming the system name is the parent folder)
    systems = defaultdict(list)
    # Track control status for each run
    control_status = {}

    for d in unique_dirs:
        sys_name = os.path.basename(os.path.dirname(d)) if os.path.dirname(d) else os.path.basename(d)
        run_name = os.path.basename(d)

        # Try to determine if this is a control system by checking summary file
        is_control = False
        summary_path = os.path.join(d, 'analysis_summary.json')
        if os.path.exists(summary_path):
            try:
                with open(summary_path, 'r') as f:
                    summary_data = json.load(f)
                    is_control = summary_data.get('IsControlSystem', False)
            except Exception as e:
                logger.warning(f"Could not read summary file {summary_path} to check control status: {e}")
                # If we can't read the file or determine control status, assume not control
                pass

        # Alternative method: check system_type.txt if summary file not available
        if not os.path.exists(summary_path):
            system_type_path = os.path.join(d, 'system_type.txt')
            if os.path.exists(system_type_path):
                try:
                    with open(system_type_path, 'r') as f:
                        lines = f.readlines()
                        for line in lines:
                            if line.strip().startswith('IS_CONTROL_SYSTEM='):
                                is_control = line.strip().split('=')[1].lower() == 'true'
                                break
                except Exception as e:
                    logger.warning(f"Could not read system_type.txt file {system_type_path} to check control status: {e}")
                    pass

        control_status[(sys_name, run_name)] = is_control
        systems[sys_name].append((run_name, d, is_control))  # Store (run_name, run_dir, is_control)

    # Sort systems alphabetically, then runs within each system
    sorted_system_names = sorted(systems.keys())

    # Define layout parameters
    left_margin = Inches(0.3)
    top_margin = Inches(1.4)
    plot_width = Inches(1.8)
    plot_height = Inches(1.3) # Reduced height to fit more rows
    h_spacing = Inches(0.15)
    v_spacing = Inches(0.15)
    max_cols = 5 # Max plots per row
    max_rows = 4 # Max plot rows per slide

    # --- Slide Generation Loop ---
    plot_types = [
        {"title_suffix": "(Raw Data)", "com_file": "COM_Stability_Plot_raw.png", "gg_file": "GG_Distance_Plot_raw.png"},
        {"title_suffix": "(Filtered Data)", "com_file": "COM_Stability_Plot.png", "gg_file": "GG_Distance_Plot.png"},
        {"title_suffix": "(Toxin Orientation)", "plot1_file": "Toxin_Orientation_Angle.png", "plot2_file": "Toxin_Channel_Contacts.png", "skip_control": True},
        {"title_suffix": "(Toxin Rotation)", "plot1_file": "Toxin_Rotation_Components.png", "plot2_file": None, "skip_control": True}, # Only one plot for rotation usually
        {"title_suffix": "(Focused Contacts)", "plot1_file": "Toxin_Channel_Residue_Contact_Map_Focused.png", "plot2_file": None, "skip_control": True},
        {"title_suffix": "(Ion Positions)", "plot1_file": "K_Ion_Combined_Plot.png", "plot2_file": None},
        {"title_suffix": "(Ion Occupancy)", "plot1_file": "K_Ion_Occupancy_Heatmap.png", "plot2_file": "K_Ion_Average_Occupancy.png"},
        {"title_suffix": "(Cavity Water)", "plot1_file": "Cavity_Water_Count_Plot.png", "plot2_file": "Cavity_Water_Residence_Hist.png"},
        # Add G1 gyration plot
        {"title_suffix": "(G1 Gyration Radius)", "plot1_file": "G1_gyration_radii.png", "plot2_file": None}
    ]

    for plot_info in plot_types:
        current_slide = None
        plots_on_current_slide = 0
        plot_idx_on_slide = 0
        plots_per_slide = max_cols * max_rows
        skip_control = plot_info.get('skip_control', False)

        for sys_name in sorted_system_names:
            runs = sorted(systems[sys_name], key=lambda x: x[0]) # Sort runs (R1, R2...)

            for i, (run_name, run_dir, is_control) in enumerate(runs):
                # Skip control systems for toxin-specific plots
                if skip_control and is_control:
                    continue

                # --- Create new slide if needed ---
                if plot_idx_on_slide % plots_per_slide == 0:
                     slide_layout = prs.slide_layouts[5] # Title only layout
                     current_slide = prs.slides.add_slide(slide_layout)
                     title = current_slide.shapes.title
                     # Try to find the first non-control run name for the title if skipping controls
                     first_run_name_for_title = run_name
                     if skip_control:
                         for rn, rd, ic in runs[i:]:
                             if not ic:
                                 first_run_name_for_title = rn
                                 break

                     title.text = f"{sys_name} - {first_run_name_for_title} ... {plot_info['title_suffix']}" # Indicate start run
                     # Reset plot index for the new slide
                     plot_idx_on_slide = 0

                # --- Calculate position ---
                row = plot_idx_on_slide // max_cols
                col = plot_idx_on_slide % max_cols
                left = left_margin + col * (plot_width + h_spacing)
                top = top_margin + row * ( (plot_height * 2) + v_spacing if plot_info.get('plot2_file') else plot_height + v_spacing)

                # --- Add plots to the current slide ---
                added_plot = False
                plot1_path = os.path.join(run_dir, plot_info.get('plot1_file', plot_info.get('com_file'))) # Handle different key names
                plot2_path = os.path.join(run_dir, plot_info.get('plot2_file', plot_info.get('gg_file'))) if plot_info.get('plot2_file') or plot_info.get('gg_file') else None

                # Add first plot (e.g., COM / Angle / Rotation / Focused Contact / Ion Combined / Water Count)
                if plot1_path and os.path.exists(plot1_path):
                    try:
                         current_slide.shapes.add_picture(plot1_path, left, top, width=plot_width, height=plot_height)
                         added_plot = True
                    except Exception as e:
                         logger.warning(f"Could not add picture {plot1_path} to PPT: {e}")
                # Add second plot if it exists (e.g., GG / Contacts / Ion Avg Occ / Water Hist)
                if plot2_path and os.path.exists(plot2_path):
                    try:
                        # Position below the first plot
                        current_slide.shapes.add_picture(plot2_path, left, top + plot_height + Inches(0.05), width=plot_width, height=plot_height)
                        added_plot = True # Count as one position even if two plots vertically
                    except Exception as e:
                         logger.warning(f"Could not add picture {plot2_path} to PPT: {e}")

                if added_plot:
                     # Add run name label with control system indicator if applicable
                     label_text = f"{sys_name}/{run_name}"
                     if is_control:
                         label_text += " (Control)" # Indicate control system
                     text_box = current_slide.shapes.add_textbox(left, top + (plot_height * 2 if plot2_path and os.path.exists(plot2_path) else plot_height) + Inches(0.05), plot_width, Inches(0.2))
                     tf = text_box.text_frame
                     tf.text = label_text
                     tf.paragraphs[0].font.size = Pt(7)
                     tf.word_wrap = False
                     plot_idx_on_slide += 1 # Increment position index only if plot(s) were added

    # --- SUMMARY TABLE SLIDE --- (Add control indicator here too)
    if com_averages:
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "COM Distance Summary (Filtered Avg)"

        # Determine table dimensions
        max_runs_per_system = 0
        for sys_name in sorted_system_names:
            max_runs_per_system = max(max_runs_per_system, len(systems[sys_name]))

        cols = 2 + max_runs_per_system # System, Avg, R1, R2...
        rows = len(sorted_system_names) + 1 # Header + one row per system

        # Adjust table size/position
        left = Inches(0.5); top = Inches(1.5)
        width = Inches(min(9.0, 1.5 * cols)) # Adjust width based on cols
        height = Inches(min(5.0, 0.3 * rows)) # Adjust height based on rows

        try:
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table

            # Header row
            table.cell(0, 0).text = "System"
            table.cell(0, 1).text = "Avg COM (Å)"
            for i in range(max_runs_per_system): table.cell(0, i + 2).text = f"R{i+1}"

            # Data rows
            row_idx = 1
            for sys_name in sorted_system_names:
                system_runs = sorted(systems[sys_name], key=lambda x: x[0]) # Get sorted runs for this system
                avg_list = []
                control_in_system = False # Flag if any run in this system is control
                for run_name, run_dir, is_control in system_runs:
                    # Look up the specific run's COM average
                    run_key = (sys_name, run_name)
                    # Find the index of this run in the original com_averages dict list if needed
                    com_vals = com_averages.get(sys_name, [])
                    if com_vals and isinstance(com_vals, list) and len(com_vals) > 0:
                        try:
                            run_idx = next((i for i, (rn, _, _) in enumerate(system_runs) if rn == run_name), None)
                            com_val = com_vals[run_idx] if run_idx is not None and run_idx < len(com_vals) else None
                            avg_list.append(com_val)
                        except Exception:
                            avg_list.append(None)
                    else:
                        avg_list.append(None)
                    
                    if is_control:
                        control_in_system = True

                # System Name and Control Indicator
                system_label = sys_name
                if control_in_system:
                    system_label += " (Control)"
                table.cell(row_idx, 0).text = system_label

                # System Average
                filtered_vals = [v for v in avg_list if v is not None and not np.isnan(v)]
                if filtered_vals:
                    system_avg = sum(filtered_vals) / len(filtered_vals)
                    table.cell(row_idx, 1).text = f"{system_avg:.1f}"
                else:
                    table.cell(row_idx, 1).text = "N/A"

                # Individual Run Values
                for i, val in enumerate(avg_list):
                    if i < max_runs_per_system: # Ensure we don't exceed table dimensions
                        if val is not None and not np.isnan(val):
                            table.cell(row_idx, i + 2).text = f"{val:.1f}"
                        else:
                            table.cell(row_idx, i + 2).text = "N/A"

                row_idx += 1

            # Format the table
            for cell in table.iter_cells():
                cell.text_frame.paragraphs[0].font.size = Pt(9)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE # Center align vertically
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER # Center align horizontally

            # Format header row as bold
            for cell in table.rows[0].cells:
                cell.text_frame.paragraphs[0].font.bold = True

        except Exception as e:
            logger.error(f"Error creating summary table: {e}", exc_info=True)
            # Fall back to simpler text box if table fails
            text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            tf = text_box.text_frame
            tf.text = "Error creating summary table. See log for details."

    # --- Save the PowerPoint ---
    ppt_path = os.path.join(os.path.dirname(unique_dirs[0]) if unique_dirs else os.getcwd(), 'MD_Analysis_Summary.pptx')
    try:
        prs.save(ppt_path)
        logger.info(f"PowerPoint saved to {ppt_path}")
    except Exception as e:
        logger.error(f"Error saving PowerPoint: {e}", exc_info=True)

    return ppt_path 