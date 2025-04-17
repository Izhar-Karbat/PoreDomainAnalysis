# reporting.py
"""
Functions for generating analysis reports (HTML for single runs, PPT for summaries).
"""

import os
import logging
import base64
import traceback
from datetime import datetime
import numpy as np
import pandas as pd
import json # <<< Add json import
import jinja2 # For HTML templating
from pptx import Presentation
from pptx.util import Inches, Pt
from collections import defaultdict
import matplotlib.pyplot as plt # Keep for plt.close if needed, though plots are loaded now
import seaborn as sns # Keep for plot style consistency settings if applied here


# Import from other modules
try:
    from utils import fig_to_base64 # Used if we were generating plots here, now used for loading
    # Plotting functions are no longer called directly here for HTML report
    # from core_analysis import plot_pore_diameter, plot_com_positions, plot_filtering_comparison, plot_kde_analysis
except ImportError as e:
    print(f"Error importing dependency modules in reporting.py: {e}")
    raise

# Get a logger for this module
logger = logging.getLogger(__name__)

# --- HTML Report Generation ---

# Define the HTML template structure here
HTML_TEMPLATE = """
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{{ run_summary.RunName }} - MD Analysis Report</title>
    <style>
        body { font-family: sans-serif; line-height: 1.6; margin: 20px; background-color: #f4f4f4; color: #333; }
        .container { max-width: 1200px; margin: auto; background: white; padding: 20px; border-radius: 8px; box-shadow: 0 0 15px rgba(0,0,0,0.1); }
        h1, h2, h3 { color: #0056b3; border-bottom: 2px solid #0056b3; padding-bottom: 5px; }
        h1 { text-align: center; }
        .section { margin-bottom: 30px; padding-bottom: 20px; border-bottom: 1px solid #eee; }
        .section:last-child { border-bottom: none; }
        .plot-container { margin: 20px 0; text-align: center; }
        .plot-container img { max-width: 95%; height: auto; border: 1px solid #ddd; padding: 5px; background: #fff; margin-bottom: 5px; }
        .stats-table { width: 100%; border-collapse: collapse; margin: 20px 0; font-size: 0.9em;}
        .stats-table th, .stats-table td { padding: 8px 10px; border: 1px solid #ddd; text-align: left; }
        .stats-table th { background-color: #e9ecef; font-weight: bold; }
        .stats-table tr:nth-child(even) { background-color: #f8f9fa; }
        .two-column { display: flex; flex-wrap: wrap; gap: 20px; }
        .column { flex: 1; min-width: 48%; } /* Adjusted min-width */
        .info-box { background-color: #e7f1ff; border-left: 5px solid #0056b3; padding: 15px; margin: 15px 0; border-radius: 4px; font-size: 0.9em;}
        .warning-box { background-color: #fff3cd; border-left: 5px solid #ffc107; padding: 15px; margin: 15px 0; border-radius: 4px; font-size: 0.9em;}
        .warning-box h4 { color: #856404; margin-top: 0; }
        .error-box { background-color: #f8d7da; border-left: 5px solid #dc3545; padding: 15px; margin: 15px 0; border-radius: 4px; font-size: 0.9em;}
        .error-box h4 { color: #721c24; margin-top: 0; }
        .control-system-banner { background-color: #17a2b8; color: white; text-align: center; padding: 10px; margin: 15px 0; border-radius: 4px; font-weight: bold; font-size: 1.1em; }
        pre { background-color: #f8f9fa; padding: 10px; border: 1px solid #ddd; border-radius: 4px; white-space: pre-wrap; word-wrap: break-word; font-size: 0.9em; max-height: 300px; overflow-y: auto;}
        .footer { text-align: center; margin-top: 30px; font-size: 0.8em; color: #666; }
        td, th { word-wrap: break-word; } /* Prevent table cell overflow */
    </style>
</head>
<body>
    <div class="container">
        <h1>MD Analysis Report</h1>
        {% if run_summary.IsControlSystem %}
        <div class="control-system-banner">CONTROL SYSTEM (NO TOXIN)</div>
        {% endif %}
        <p class="footer">Generated: {{ generation_timestamp }} | Analysis Version: {{ run_summary.AnalysisScriptVersion }}</p>

        <div class="section">
            <h2>Run Information</h2>
            <table class="stats-table">
                <tr><th>System</th><td>{{ run_summary.SystemName }}</td></tr>
                <tr><th>Run</th><td>{{ run_summary.RunName }}</td></tr>
                <tr><th>Path</th><td>{{ run_summary.RunPath }}</td></tr>
                <tr><th>System Type</th><td>{% if run_summary.IsControlSystem %}Control (No Toxin){% else %}Toxin-Channel Complex{% endif %}</td></tr>
                <tr><th>Analysis Status</th><td>{{ run_summary.AnalysisStatus }}</td></tr>
                <tr><th>Analysis Timestamp</th><td>{{ run_summary.AnalysisTimestamp }}</td></tr>
            </table>
        </div>

        <div class="section">
            <h2>G-G Distance Analysis (Pore Diameter Proxy)</h2>
            <table class="stats-table">
              <thead><tr><th>Metric</th><th>Filtered A:C</th><th>Filtered B:D</th></tr></thead>
              <tbody>
                <tr><td>Mean (Å)</td><td>{{ "%.3f"|format(run_summary.GG_AC_Mean_Filt) if run_summary.GG_AC_Mean_Filt is not none else 'N/A' }}</td><td>{{ "%.3f"|format(run_summary.GG_BD_Mean_Filt) if run_summary.GG_BD_Mean_Filt is not none else 'N/A' }}</td></tr>
                <tr><td>Std Dev (Å)</td><td>{{ "%.3f"|format(run_summary.GG_AC_Std_Filt) if run_summary.GG_AC_Std_Filt is not none else 'N/A' }}</td><td>{{ "%.3f"|format(run_summary.GG_BD_Std_Filt) if run_summary.GG_BD_Std_Filt is not none else 'N/A' }}</td></tr>
                <tr><td>Min (Å)</td><td>{{ "%.3f"|format(run_summary.GG_AC_Min_Filt) if run_summary.GG_AC_Min_Filt is not none else 'N/A' }}</td><td>{{ "%.3f"|format(run_summary.GG_BD_Min_Filt) if run_summary.GG_BD_Min_Filt is not none else 'N/A' }}</td></tr>
              </tbody>
            </table>
            <div class="two-column">
                <div class="column plot-container">
                    <h3>Raw G-G Distance</h3>
                    {% if img_data.GG_Distance_Plot_raw %} <img src="data:image/png;base64,{{ img_data.GG_Distance_Plot_raw }}" alt="Raw G-G Distance Plot"> {% else %} <p><i>Plot not available.</i></p> {% endif %}
                </div>
                <div class="column plot-container">
                    <h3>Filtered G-G Distance</h3>
                    {% if img_data.GG_Distance_Plot %} <img src="data:image/png;base64,{{ img_data.GG_Distance_Plot }}" alt="Filtered G-G Distance Plot"> {% else %} <p><i>Plot not available.</i></p> {% endif %}
                </div>
            </div>
            {% if img_data.G_G_Distance_AC_Comparison %}
            <div class="plot-container">
                <h3>A:C Distance Filtering Detail</h3>
                <img src="data:image/png;base64,{{ img_data.G_G_Distance_AC_Comparison }}" alt="A:C Filtering Comparison">
            </div>
            {% endif %}
            {% if img_data.G_G_Distance_BD_Comparison %}
            <div class="plot-container">
                <h3>B:D Distance Filtering Detail</h3>
                <img src="data:image/png;base64,{{ img_data.G_G_Distance_BD_Comparison }}" alt="B:D Filtering Comparison">
            </div>
            {% endif %}
        </div>

        {% if not run_summary.IsControlSystem and run_summary.COM_Mean_Filt is not none %} {# Only show for non-control systems with COM data #}
        <div class="section">
            <h2>COM Distance Analysis (Toxin-Channel Stability)</h2>
            <div class="info-box">
                Filter Type: {{ run_summary.COM_Filter_Type | default('N/A') }} |
                Quality Issue Detected: {{ 'Yes' if run_summary.COM_Filter_QualityIssue else 'No' }}
            </div>
            <table class="stats-table">
              <thead><tr><th>Metric</th><th>Filtered COM</th></tr></thead>
              <tbody>
                <tr><td>Mean (Å)</td><td>{{ "%.3f"|format(run_summary.COM_Mean_Filt) }}</td></tr>
                <tr><td>Std Dev (Å)</td><td>{{ "%.3f"|format(run_summary.COM_Std_Filt) }}</td></tr>
                <tr><td>Min (Å)</td><td>{{ "%.3f"|format(run_summary.COM_Min_Filt) }}</td></tr>
                <tr><td>Max (Å)</td><td>{{ "%.3f"|format(run_summary.COM_Max_Filt) }}</td></tr>
              </tbody>
            </table>
            <div class="two-column">
                <div class="column plot-container">
                    <h3>Raw COM Distance</h3>
                    {% if img_data.COM_Stability_Plot_raw %} <img src="data:image/png;base64,{{ img_data.COM_Stability_Plot_raw }}" alt="Raw COM Stability Plot"> {% else %} <p><i>Plot not available.</i></p> {% endif %}
                </div>
                <div class="column plot-container">
                    <h3>Filtered COM Distance</h3>
                     {% if img_data.COM_Stability_Plot %} <img src="data:image/png;base64,{{ img_data.COM_Stability_Plot }}" alt="Filtered COM Stability Plot"> {% else %} <p><i>Plot not available.</i></p> {% endif %}
                </div>
            </div>
            {% if img_data.COM_Stability_Comparison %}
            <div class="plot-container">
                <h3>COM Distance Filtering Detail</h3>
                <img src="data:image/png;base64,{{ img_data.COM_Stability_Comparison }}" alt="COM Filtering Comparison">
            </div>
            {% endif %}
            {% if img_data.COM_Stability_KDE_Analysis %}
            <div class="plot-container">
                <h3>COM Level Analysis (KDE on Raw Data)</h3>
                <img src="data:image/png;base64,{{ img_data.COM_Stability_KDE_Analysis }}" alt="COM KDE Analysis">
            </div>
            {% endif %}
        </div>
        {% elif run_summary.IsControlSystem %}
        <div class="section">
            <h2>COM Distance Analysis</h2>
            <div class="info-box">
                This is a control system without toxin. COM distance analysis is not applicable.
            </div>
        </div>
        {% else %}
        <div class="section"><p><i>COM distance analysis not performed or data unavailable.</i></p></div>
        {% endif %}

        {% if not run_summary.IsControlSystem and run_summary.COM_Mean_Filt is not none %} {# Only show for non-control systems #}
        <div class="section">
            <h2>Toxin Orientation & Interface Analysis</h2>
             <table class="stats-table">
               <thead><tr><th>Metric</th><th>Value</th></tr></thead>
               <tbody>
                 <tr><td>Mean Orientation Angle (°)</td><td>{{ "%.2f"|format(run_summary.Orient_Angle_Mean) if run_summary.Orient_Angle_Mean is not none else 'N/A' }}</td></tr>
                 <tr><td>Std Dev Orientation Angle (°)</td><td>{{ "%.2f"|format(run_summary.Orient_Angle_Std) if run_summary.Orient_Angle_Std is not none else 'N/A' }}</td></tr>
                 <tr><td>Mean Atom Contacts</td><td>{{ "%.1f"|format(run_summary.Orient_Contacts_Mean) if run_summary.Orient_Contacts_Mean is not none else 'N/A' }}</td></tr>
                 <tr><td>Std Dev Atom Contacts</td><td>{{ "%.1f"|format(run_summary.Orient_Contacts_Std) if run_summary.Orient_Contacts_Std is not none else 'N/A' }}</td></tr>
               </tbody>
             </table>
            <div class="two-column">
                <div class="column plot-container">
                    <h3>Orientation Angle</h3>
                    {% if img_data.Toxin_Orientation_Angle %} <img src="data:image/png;base64,{{ img_data.Toxin_Orientation_Angle }}" alt="Toxin Orientation Angle"> {% else %} <p><i>Plot not available.</i></p> {% endif %}
                </div>
                <div class="column plot-container">
                    <h3>Rotation Components</h3>
                    {% if img_data.Toxin_Rotation_Components %} <img src="data:image/png;base64,{{ img_data.Toxin_Rotation_Components }}" alt="Toxin Rotation Components"> {% else %} <p><i>Plot not available.</i></p> {% endif %}
                </div>
            </div>
             <div class="plot-container">
                 <h3>Total Atom Contacts</h3>
                 {% if img_data.Toxin_Channel_Contacts %} <img src="data:image/png;base64,{{ img_data.Toxin_Channel_Contacts }}" alt="Toxin-Channel Contacts"> {% else %} <p><i>Plot not available.</i></p> {% endif %}
             </div>
             <div class="plot-container">
                 <h3>Focused Residue Contact Map</h3>
                 {% if img_data.Toxin_Channel_Residue_Contact_Map_Focused %} <img src="data:image/png;base64,{{ img_data.Toxin_Channel_Residue_Contact_Map_Focused }}" alt="Focused Contact Map"> {% else %} <p><i>Plot not available.</i></p> {% endif %}
             </div>
             <div class="plot-container">
                 <h3>Full Residue Contact Map</h3>
                 {% if img_data.Toxin_Channel_Residue_Contact_Map_Full %} <img src="data:image/png;base64,{{ img_data.Toxin_Channel_Residue_Contact_Map_Full }}" alt="Full Contact Map"> {% else %} <p><i>Plot not available.</i></p> {% endif %}
             </div>
        </div>
        {% elif run_summary.IsControlSystem %}
        <div class="section">
            <h2>Toxin Orientation & Interface Analysis</h2>
            <div class="info-box">
                This is a control system without toxin. Orientation and interface analysis is not applicable.
            </div>
        </div>
        {% endif %}

        <div class="section">
            <h2>K+ Ion Analysis (Selectivity Filter)</h2>
            {% if run_summary.Ion_Count > 0 %} {# Check if ions were tracked #}
                 <div class="info-box">
                    Tracked {{ run_summary.Ion_Count }} unique K+ ions near the filter.<br>
                    Coordinates relative to G1 Cα Z-plane = 0 Å.
                 </div>
                 {% if binding_site_data %}
                 <div class="info-box">
                     <h3>Binding Site Positions (G1 Cα = 0 Å)</h3>
                     <pre>{{ binding_site_data }}</pre>
                 </div>
                 {% endif %}

                 {% if img_data.binding_sites_g1_centric_visualization %}
                 <div class="plot-container">
                     <h3>Binding Site Visualization</h3>
                     <img src="data:image/png;base64,{{ img_data.binding_sites_g1_centric_visualization }}" alt="Binding Sites Visualization">
                 </div>
                 {% endif %}

                 {% if img_data.K_Ion_Combined_Plot %}
                 <div class="plot-container">
                    <h3>Ion Positions & Density</h3>
                    <img src="data:image/png;base64,{{ img_data.K_Ion_Combined_Plot }}" alt="K+ Ion Positions and Density">
                 </div>
                 {% endif %}

                 {% if img_data.K_Ion_Occupancy_Heatmap %}
                 <div class="plot-container">
                    <h3>Site Occupancy Heatmap</h3>
                    <img src="data:image/png;base64,{{ img_data.K_Ion_Occupancy_Heatmap }}" alt="K+ Ion Occupancy Heatmap">
                 </div>
                 {% endif %}

                 {% if img_data.K_Ion_Average_Occupancy %}
                 <div class="plot-container">
                    <h3>Average Site Occupancy</h3>
                    <img src="data:image/png;base64,{{ img_data.K_Ion_Average_Occupancy }}" alt="Average K+ Ion Occupancy">
                 </div>
                 {% endif %}

                 <h3>Site Occupancy Statistics</h3>
                 <table class="stats-table">
                     <thead><tr><th>Site</th><th>Mean Occ.</th><th>Max Occ.</th><th>% Time Occ. (>0)</th></tr></thead>
                     <tbody>
                     {% for site in ['S0', 'S1', 'S2', 'S3', 'S4', 'Cavity'] %}
                         <tr>
                             <td>{{ site }}</td>
                             <td>{{ "%.3f"|format(run_summary['Ion_AvgOcc_'~site]) if run_summary['Ion_AvgOcc_'~site] is not none else 'N/A' }}</td>
                             <td>{{ "%d"|format(run_summary['Ion_MaxOcc_'~site]) if run_summary['Ion_MaxOcc_'~site] is not none else 'N/A' }}</td>
                             <td>{{ "%.1f"|format(run_summary['Ion_PctTimeOcc_'~site]) if run_summary['Ion_PctTimeOcc_'~site] is not none else 'N/A' }}%</td>
                         </tr>
                     {% endfor %}
                     </tbody>
                 </table>
            {% else %}
                <p><i>K+ ion analysis not performed or no ions tracked near filter.</i></p>
            {% endif %}
        </div>

        <div class="section">
            <h2>Cavity Water Analysis</h2>
             {% if run_summary.CavityWater_MeanOcc is not none %} {# Check if water stats exist #}
                 <table class="stats-table">
                    <thead><tr><th>Metric</th><th>Value</th></tr></thead>
                    <tbody>
                        <tr><td>Mean Occupancy</td><td>{{ "%.2f"|format(run_summary.CavityWater_MeanOcc) }}</td></tr>
                        <tr><td>Std Dev Occupancy</td><td>{{ "%.2f"|format(run_summary.CavityWater_StdOcc) }}</td></tr>
                        <tr><td>Avg Residence Time (ns)</td><td>{{ "%.3f"|format(run_summary.CavityWater_AvgResidenceTime_ns) }}</td></tr>
                        <tr><td>Total Confirmed Exit Events</td><td>{{ "%d"|format(run_summary.CavityWater_TotalExitEvents) }}</td></tr>
                        <tr><td>Exchange Rate (/ns)</td><td>{{ "%.4f"|format(run_summary.CavityWater_ExchangeRatePerNs) }}</td></tr>
                    </tbody>
                 </table>
                 <div class="two-column">
                    <div class="column plot-container">
                        <h3>Cavity Water Count</h3>
                         {% if img_data.Cavity_Water_Count_Plot %} <img src="data:image/png;base64,{{ img_data.Cavity_Water_Count_Plot }}" alt="Cavity Water Count Plot"> {% else %} <p><i>Plot not available.</i></p> {% endif %}
                    </div>
                    <div class="column plot-container">
                        <h3>Residence Time Distribution</h3>
                        {% if img_data.Cavity_Water_Residence_Hist %} <img src="data:image/png;base64,{{ img_data.Cavity_Water_Residence_Hist }}" alt="Cavity Water Residence Histogram"> {% else %} <p><i>Plot not available.</i></p> {% endif %}
                    </div>
                 </div>
             {% else %}
                 <p><i>Cavity water analysis not performed or data unavailable.</i></p>
             {% endif %}
        </div>

        {# <<< ADDED SECTION for Carbonyl Gyration >>> #}
        <div class="section">
            <h2>Carbonyl Gyration Analysis (Selectivity Filter G1)</h2>
             {% if run_summary.Gyration_G1_Mean is not none %} {# Check if gyration analysis was performed #}
                 <div class="info-box">
                    Gyration radius (ρ) measures the distance between the G1 carbonyl oxygen atoms and the pore center.
                    Changes in this radius can indicate carbonyl flipping events that may affect ion permeation.
                 </div>
                 <table class="stats-table">
                    <thead><tr><th>Metric</th><th>G1 Glycine</th></tr></thead>
                    <tbody>
                        <tr><td>Mean Gyration Radius (Å)</td>
                            <td>{{ "%.3f"|format(run_summary.Gyration_G1_Mean) }}</td>
                        </tr>
                        <tr><td>Std Dev (Å)</td>
                            <td>{{ "%.3f"|format(run_summary.Gyration_G1_Std) }}</td>
                        </tr>
                        <tr><td>Detected Flips</td>
                            <td>{{ run_summary.Gyration_Flips }}</td>
                        </tr>
                        <tr><td>Maximum Change (Å)</td>
                            <td>{{ "%.3f"|format(run_summary.Gyration_MaxChange) if run_summary.Gyration_MaxChange is not none else 'N/A' }}</td>
                        </tr>
                    </tbody>
                 </table>
                 <div class="plot-container">
                     <h3>G1 Gyration Radii</h3>
                     {% if img_data.G1_gyration_radii %} <img src="data:image/png;base64,{{ img_data.G1_gyration_radii }}" alt="G1 Gyration Radii"> {% else %} <p><i>Plot not available.</i></p> {% endif %}
                 </div>
                 <div class="info-box">
                    <p><strong>Interpretation:</strong> Carbonyl flipping is indicated by sudden changes in gyration radius.
                    {% if run_summary.Gyration_Flips > 5 %}
                    A high number of flips ({{ run_summary.Gyration_Flips }}) suggests significant instability in the selectivity filter,
                    which may disrupt ion coordination and permeation.
                    {% elif run_summary.Gyration_Flips > 0 %}
                    {{ run_summary.Gyration_Flips }} flip(s) were detected, which may temporarily affect ion coordination.
                    {% else %}
                    No significant flips were detected, suggesting a stable selectivity filter conformation.
                    {% endif %}
                    </p>
                 </div>
             {% else %}
                 <p><i>Carbonyl gyration analysis not performed or data unavailable.</i></p>
             {% endif %}
        </div>

        <div class="footer">End of Report</div>
    </div>
</body>
</html>
"""


def _load_image_base64(image_path):
    """Loads an image file and returns its base64 encoded string."""
    if image_path and os.path.exists(image_path):
        try:
            with open(image_path, 'rb') as f:
                return base64.b64encode(f.read()).decode('utf-8')
        except Exception as e:
            logger.warning(f"Error loading image {image_path}: {e}")
    return None

def generate_html_report(run_dir, run_summary):
    """
    Generates a comprehensive HTML report for a single run.

    Relies on existing plot files and the summary dictionary.

    Args:
        run_dir (str): Path to the specific run directory.
        run_summary (dict): Dictionary containing summary statistics for the run,
                            loaded from 'analysis_summary.json'.

    Returns:
        str | None: Path to the generated HTML report, or None on failure.
    """
    logger.info(f"Generating HTML report for {run_summary.get('RunName', 'Unknown Run')}...")

    # --- Load Base64 Encoded Images ---
    img_data = {}
    plot_files = {
        # Core distance plots (now in core_analysis/)
        "GG_Distance_Plot_raw": "core_analysis/GG_Distance_Plot_raw.png",
        "GG_Distance_Plot": "core_analysis/GG_Distance_Plot.png",
        "G_G_Distance_AC_Comparison": "core_analysis/G_G_Distance_AC_Comparison.png",
        "G_G_Distance_BD_Comparison": "core_analysis/G_G_Distance_BD_Comparison.png",
        "COM_Stability_Plot_raw": "core_analysis/COM_Stability_Plot_raw.png",
        "COM_Stability_Plot": "core_analysis/COM_Stability_Plot.png",
        "COM_Stability_Comparison": "core_analysis/COM_Stability_Comparison.png",
        "COM_Stability_KDE_Analysis": "core_analysis/COM_Stability_KDE_Analysis.png",
        # Orientation/Contact plots (now in orientation_contacts/)
        "Toxin_Orientation_Angle": "orientation_contacts/Toxin_Orientation_Angle.png",
        "Toxin_Rotation_Components": "orientation_contacts/Toxin_Rotation_Components.png",
        "Toxin_Channel_Contacts": "orientation_contacts/Toxin_Channel_Contacts.png",
        "Toxin_Channel_Residue_Contact_Map_Full": "orientation_contacts/Toxin_Channel_Residue_Contact_Map_Full.png",
        "Toxin_Channel_Residue_Contact_Map_Focused": "orientation_contacts/Toxin_Channel_Residue_Contact_Map_Focused.png",
        # Ion plots (now in ion_analysis/)
        "binding_sites_g1_centric_visualization": "ion_analysis/binding_sites_g1_centric_visualization.png",
        "K_Ion_Combined_Plot": "ion_analysis/K_Ion_Combined_Plot.png",
        "K_Ion_Occupancy_Heatmap": "ion_analysis/K_Ion_Occupancy_Heatmap.png",
        "K_Ion_Average_Occupancy": "ion_analysis/K_Ion_Average_Occupancy.png",
        # Water plots (now in water_analysis/)
        "Cavity_Water_Count_Plot": "water_analysis/Cavity_Water_Count_Plot.png",
        "Cavity_Water_Residence_Hist": "water_analysis/Cavity_Water_Residence_Hist.png",
        # Gyration analysis plots (now in gyration_analysis/)
        "G1_gyration_radii": "gyration_analysis/G1_gyration_radii.png",
    }

    logger.debug("Loading images for HTML report...")
    for key, filename in plot_files.items():
        img_path = os.path.join(run_dir, filename)
        img_data[key] = _load_image_base64(img_path)
        if img_data[key] is None:
            logger.debug(f"Plot file not found or load failed: {img_path}")

    # --- Load Binding Site Data ---
    binding_site_data = None
    binding_site_path = os.path.join(run_dir, 'binding_site_positions_g1_centric.txt')
    if os.path.exists(binding_site_path):
        try:
            with open(binding_site_path, 'r') as f:
                binding_site_data = f.read()
            logger.debug("Loaded binding site position data.")
        except Exception as e:
            logger.warning(f"Error reading binding site data from {binding_site_path}: {e}")

    # --- Render HTML ---
    try:
        logger.debug("Rendering HTML template...")
        env = jinja2.Environment(loader=jinja2.BaseLoader(), autoescape=jinja2.select_autoescape(['html', 'xml']))
        template = env.from_string(HTML_TEMPLATE)

        # Prepare context for rendering
        render_context = {
            'run_summary': run_summary, # Pass the loaded summary dictionary
            'img_data': img_data,
            'binding_site_data': binding_site_data,
            'generation_timestamp': datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        }

        rendered_html = template.render(render_context)
        logger.debug("HTML rendering complete.")

        # Save the report
        report_filename = f"{run_summary.get('RunName', 'report')}_analysis_report.html"
        report_path = os.path.join(run_dir, report_filename)
        with open(report_path, 'w', encoding='utf-8') as f:
            f.write(rendered_html)

        logger.info(f"Generated HTML report at: {report_path}")
        return report_path

    except jinja2.exceptions.TemplateError as e:
        logger.error(f"Error rendering Jinja2 template: {e}", exc_info=True)
        error_report_path = os.path.join(run_dir, f"{run_summary.get('RunName', 'report')}_report_TEMPLATE_ERROR.html")
        with open(error_report_path, 'w') as f: f.write(f"<html><body><h1>Template Error</h1><pre>{traceback.format_exc()}</pre></body></html>")
        return None
    except Exception as e:
        logger.error(f"Unexpected error generating HTML report: {e}", exc_info=True)
        error_report_path = os.path.join(run_dir, f"{run_summary.get('RunName', 'report')}_report_ERROR.html")
        with open(error_report_path, 'w') as f: f.write(f"<html><body><h1>Report Generation Error</h1><pre>{traceback.format_exc()}</pre></body></html>")
        return None

# --- PowerPoint Report Generation ---

def Create_PPT(unique_dirs, com_averages):
    """
    Creates a PowerPoint presentation summarizing results from multiple runs.
    Includes raw/filtered distance plots, orientation plots, and a COM summary table.
    Clearly identifies control systems vs toxin-channel systems.

    Args:
        unique_dirs (list): List of paths to the run directories to include.
        com_averages (dict): Dictionary mapping system names to lists of filtered
                             average COM distances for the runs within that system.
    """
    logger.info(f"Creating PowerPoint summary for {len(unique_dirs)} runs...")
    prs = Presentation()

    # Group runs by system (assuming the system name is the parent folder)
    systems = defaultdict(list)
    # Track control status for each run
    control_status = {}

    for d in unique_dirs:
        sys_name = os.path.basename(os.path.dirname(d)) if os.path.dirname(d) else os.path.basename(d)
        run_name = os.path.basename(d)

        # Try to determine if this is a control system by checking summary file
        is_control = False
        summary_path = os.path.join(d, 'analysis_summary.json')
        if os.path.exists(summary_path):
            try:
                with open(summary_path, 'r') as f:
                    summary_data = json.load(f)
                    is_control = summary_data.get('IsControlSystem', False)
            except Exception as e:
                logger.warning(f"Could not read summary file {summary_path} to check control status: {e}")
                # If we can't read the file or determine control status, assume not control
                pass

        # Alternative method: check system_type.txt if summary file not available
        if not os.path.exists(summary_path):
            system_type_path = os.path.join(d, 'system_type.txt')
            if os.path.exists(system_type_path):
                try:
                    with open(system_type_path, 'r') as f:
                        lines = f.readlines()
                        for line in lines:
                            if line.strip().startswith('IS_CONTROL_SYSTEM='):
                                is_control = line.strip().split('=')[1].lower() == 'true'
                                break
                except Exception as e:
                    logger.warning(f"Could not read system_type.txt file {system_type_path} to check control status: {e}")
                    pass

        control_status[(sys_name, run_name)] = is_control
        systems[sys_name].append((run_name, d, is_control))  # Store (run_name, run_dir, is_control)

    # Sort systems alphabetically, then runs within each system
    sorted_system_names = sorted(systems.keys())

    # Define layout parameters
    left_margin = Inches(0.3)
    top_margin = Inches(1.4)
    plot_width = Inches(1.8)
    plot_height = Inches(1.3) # Reduced height to fit more rows
    h_spacing = Inches(0.15)
    v_spacing = Inches(0.15)
    max_cols = 5 # Max plots per row
    max_rows = 4 # Max plot rows per slide

    # --- Slide Generation Loop ---
    plot_types = [
        {"title_suffix": "(Raw Data)", "com_file": "COM_Stability_Plot_raw.png", "gg_file": "GG_Distance_Plot_raw.png"},
        {"title_suffix": "(Filtered Data)", "com_file": "COM_Stability_Plot.png", "gg_file": "GG_Distance_Plot.png"},
        {"title_suffix": "(Toxin Orientation)", "plot1_file": "Toxin_Orientation_Angle.png", "plot2_file": "Toxin_Channel_Contacts.png", "skip_control": True},
        {"title_suffix": "(Toxin Rotation)", "plot1_file": "Toxin_Rotation_Components.png", "plot2_file": None, "skip_control": True}, # Only one plot for rotation usually
        {"title_suffix": "(Focused Contacts)", "plot1_file": "Toxin_Channel_Residue_Contact_Map_Focused.png", "plot2_file": None, "skip_control": True},
        {"title_suffix": "(Ion Positions)", "plot1_file": "K_Ion_Combined_Plot.png", "plot2_file": None},
        {"title_suffix": "(Ion Occupancy)", "plot1_file": "K_Ion_Occupancy_Heatmap.png", "plot2_file": "K_Ion_Average_Occupancy.png"},
        {"title_suffix": "(Cavity Water)", "plot1_file": "Cavity_Water_Count_Plot.png", "plot2_file": "Cavity_Water_Residence_Hist.png"},
        # Add G1 gyration plot
        {"title_suffix": "(G1 Gyration Radius)", "plot1_file": "G1_gyration_radii.png", "plot2_file": None}
    ]

    for plot_info in plot_types:
        current_slide = None
        plots_on_current_slide = 0
        plot_idx_on_slide = 0
        plots_per_slide = max_cols * max_rows
        skip_control = plot_info.get('skip_control', False)

        for sys_name in sorted_system_names:
            runs = sorted(systems[sys_name], key=lambda x: x[0]) # Sort runs (R1, R2...)

            for i, (run_name, run_dir, is_control) in enumerate(runs):
                # Skip control systems for toxin-specific plots
                if skip_control and is_control:
                    continue

                # --- Create new slide if needed ---
                if plot_idx_on_slide % plots_per_slide == 0:
                     slide_layout = prs.slide_layouts[5] # Title only layout
                     current_slide = prs.slides.add_slide(slide_layout)
                     title = current_slide.shapes.title
                     # Try to find the first non-control run name for the title if skipping controls
                     first_run_name_for_title = run_name
                     if skip_control:
                         for rn, rd, ic in runs[i:]:
                             if not ic:
                                 first_run_name_for_title = rn
                                 break

                     title.text = f"{sys_name} - {first_run_name_for_title} ... {plot_info['title_suffix']}" # Indicate start run
                     # Reset plot index for the new slide
                     plot_idx_on_slide = 0

                # --- Calculate position ---
                row = plot_idx_on_slide // max_cols
                col = plot_idx_on_slide % max_cols
                left = left_margin + col * (plot_width + h_spacing)
                top = top_margin + row * ( (plot_height * 2) + v_spacing if plot_info.get('plot2_file') else plot_height + v_spacing)

                # --- Add plots to the current slide ---
                added_plot = False
                plot1_path = os.path.join(run_dir, plot_info.get('plot1_file', plot_info.get('com_file'))) # Handle different key names
                plot2_path = os.path.join(run_dir, plot_info.get('plot2_file', plot_info.get('gg_file'))) if plot_info.get('plot2_file') or plot_info.get('gg_file') else None

                # Add first plot (e.g., COM / Angle / Rotation / Focused Contact / Ion Combined / Water Count)
                if plot1_path and os.path.exists(plot1_path):
                    try:
                         current_slide.shapes.add_picture(plot1_path, left, top, width=plot_width, height=plot_height)
                         added_plot = True
                    except Exception as e:
                         logger.warning(f"Could not add picture {plot1_path} to PPT: {e}")
                # Add second plot if it exists (e.g., GG / Contacts / Ion Avg Occ / Water Hist)
                if plot2_path and os.path.exists(plot2_path):
                    try:
                        # Position below the first plot
                        current_slide.shapes.add_picture(plot2_path, left, top + plot_height + Inches(0.05), width=plot_width, height=plot_height)
                        added_plot = True # Count as one position even if two plots vertically
                    except Exception as e:
                         logger.warning(f"Could not add picture {plot2_path} to PPT: {e}")

                if added_plot:
                     # Add run name label with control system indicator if applicable
                     label_text = f"{sys_name}/{run_name}"
                     if is_control:
                         label_text += " (Control)" # Indicate control system
                     text_box = current_slide.shapes.add_textbox(left, top + (plot_height * 2 if plot2_path and os.path.exists(plot2_path) else plot_height) + Inches(0.05), plot_width, Inches(0.2))
                     tf = text_box.text_frame
                     tf.text = label_text
                     tf.paragraphs[0].font.size = Pt(7)
                     tf.word_wrap = False
                     plot_idx_on_slide += 1 # Increment position index only if plot(s) were added

    # --- SUMMARY TABLE SLIDE --- (Add control indicator here too)
    if com_averages:
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "COM Distance Summary (Filtered Avg)"

        # Determine table dimensions
        max_runs_per_system = 0
        for sys_name in sorted_system_names:
            max_runs_per_system = max(max_runs_per_system, len(systems[sys_name]))

        cols = 2 + max_runs_per_system # System, Avg, R1, R2...
        rows = len(sorted_system_names) + 1 # Header + one row per system

        # Adjust table size/position
        left = Inches(0.5); top = Inches(1.5)
        width = Inches(min(9.0, 1.5 * cols)) # Adjust width based on cols
        height = Inches(min(5.0, 0.3 * rows)) # Adjust height based on rows

        try:
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table

            # Header row
            table.cell(0, 0).text = "System"
            table.cell(0, 1).text = "Avg COM (Å)"
            for i in range(max_runs_per_system): table.cell(0, i + 2).text = f"R{i+1}"

            # Data rows
            row_idx = 1
            for sys_name in sorted_system_names:
                system_runs = sorted(systems[sys_name], key=lambda x: x[0]) # Get sorted runs for this system
                avg_list = []
                control_in_system = False # Flag if any run in this system is control
                for run_name, run_dir, is_control in system_runs:
                    # Look up the specific run's COM average
                    run_key = (sys_name, run_name)
                    # Find the index of this run in the original com_averages dict list if needed
                    # This assumes com_averages[sys_name] is ordered same as system_runs - simplified here
                    com_val = com_averages.get(sys_name, {}).get(run_name, None)
                    avg_list.append(com_val)
                    if is_control:
                        control_in_system = True

                # System Name and Control Indicator
                system_label = sys_name
                if control_in_system:
                    # Could check if *all* are control vs mixed
                    all_control = all(run[2] for run in system_runs)
                    if all_control:
                         system_label += " (Control System)"
                    else:
                         system_label += " (Mixed Runs)" # Indicate mix of control/non-control
                table.cell(row_idx, 0).text = system_label

                # Calculate Overall Average COM (excluding None/NaN and potentially excluding control runs)
                valid_coms = [val for val, run_info in zip(avg_list, system_runs) if val is not None and not np.isnan(val) and not run_info[2]] # Exclude controls from avg
                overall_avg = np.mean(valid_coms) if valid_coms else np.nan
                table.cell(row_idx, 1).text = f"{overall_avg:.3f}" if not np.isnan(overall_avg) else "N/A (or Control Only)"

                # Fill in individual run COM values
                for j in range(max_runs_per_system):
                    if j < len(system_runs):
                        run_name, run_dir, is_control = system_runs[j]
                        com_val = avg_list[j]
                        cell_text = f"{com_val:.3f}" if com_val is not None and not np.isnan(com_val) else ("N/A" if is_control else "")
                        if is_control:
                            cell_text += " (C)" # Indicator for control run
                        table.cell(row_idx, j + 2).text = cell_text
                    else:
                        table.cell(row_idx, j + 2).text = ""
                row_idx += 1

             # --- Optional: Apply basic table styling ---
            for col_idx in range(cols): table.columns[col_idx].width = Inches(1.1) # Adjust column width
            for row in table.rows: row.height = Inches(0.25)
            # Add more styling as needed (font size, alignment...)
            # e.g., Set font size for the whole table
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(8)

        except Exception as e:
            logger.error(f"Failed to create summary table in PowerPoint: {e}", exc_info=True)
            # Add a text box indicating the error
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = f"Error creating summary table:\n{e}"

    # --- Save Presentation ---
    out_ppt = "MD_Analysis_Summary.pptx" # Consistent name
    try:
        prs.save(out_ppt)
        logger.info(f"PowerPoint presentation saved as '{out_ppt}'.")
    except PermissionError as e:
        logger.error(f"Failed to save PowerPoint '{out_ppt}'. Check if the file is open or permissions are correct: {e}")
    except Exception as e:
        logger.error(f"Failed to save PowerPoint presentation '{out_ppt}': {e}", exc_info=True)
